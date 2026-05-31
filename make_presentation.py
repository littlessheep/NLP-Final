from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(".")
OUT = ROOT / "NLP_Final_Project_Presentation.pptx"

IMG = {
    "top_words": ROOT / "output/preprocessing_stats/top_words_bar.png",
    "doc_len": ROOT / "output/preprocessing_stats/doc_length_distribution.png",
    "k_score": ROOT / "output/lda_k_selection/lda_k_composite_score.png",
    "k_npmi": ROOT / "output/lda_k_selection/lda_k_npmi_coherence.png",
    "topic_words": ROOT / "output/lda_k15/topic_visualizations/topic_top_words_facets.png",
    "topic_prev": ROOT / "output/lda_k15/topic_visualizations/topic_prevalence_bar.png",
    "dominant": ROOT / "output/lda_k15/topic_visualizations/dominant_topic_distribution.png",
    "corr": ROOT / "output/lda_k15/topic_visualizations/topic_correlation_heatmap.png",
    "dendro": ROOT / "output/lda_k15/topic_visualizations/topic_clustering_dendrogram.png",
    "doc_topic": ROOT / "output/lda_k15/topic_visualizations/document_topic_heatmap_sample.png",
}


TITLE = RGBColor(26, 54, 93)
TEXT = RGBColor(35, 35, 35)
MUTED = RGBColor(90, 100, 115)
ACCENT = RGBColor(47, 111, 167)
LIGHT = RGBColor(239, 244, 249)


def set_font(run, size=20, bold=False, color=TEXT):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=20, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.28, 12.2, 0.5, title, size=25, bold=True, color=TITLE)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(0.92), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.58, 1.02, 12.1, 0.35, subtitle, size=12, color=MUTED)


def add_bullets(slide, x, y, w, h, bullets, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(7)
    return box


def add_image(slide, path, x, y, w=None, h=None):
    if not path.exists():
        add_textbox(slide, x, y, w or 5, h or 1, f"缺少图片：{path}", size=14, color=RGBColor(160, 50, 50))
        return None
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_table(slide, x, y, w, h, rows, cols, data, font_size=11):
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Microsoft YaHei"
            para.font.size = Pt(font_size)
            para.font.color.rgb = TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
                para.font.bold = True
    return table


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def make_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = blank(prs)
    add_textbox(slide, 0.8, 1.35, 11.7, 0.9, "英文摘要语料的主题建模分析", size=36, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.6, 2.35, 10.2, 0.6, "NLP Final Project：数据清洗、文本预处理、LDA 主题模型与可视化", size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.15, 4.45, 9.0, 0.5, "数据规模：12195 篇英文摘要｜最终模型：LDA, K=15", size=18, color=TEXT, align=PP_ALIGN.CENTER)

    # 2
    slide = blank(prs)
    add_title(slide, "研究目标与整体流程")
    add_bullets(slide, 0.8, 1.45, 5.7, 4.9, [
        "目标：从英文摘要语料中识别主要研究主题",
        "方法：先进行文本清洗与矩阵构建，再使用 LDA 进行主题建模",
        "输出：主题关键词、代表文档、主题占比、主题相关性和聚类图",
        "最终用于支持项目报告中的 Method、Results 和 Discussion",
    ], size=18)
    add_textbox(slide, 7.0, 1.55, 5.4, 0.45, "项目流程", size=20, bold=True, color=TITLE)
    add_bullets(slide, 7.0, 2.15, 5.5, 4.0, [
        "1. 原始摘要清洗",
        "2. 分词、停用词过滤、词形还原",
        "3. 构建 Count Matrix 和 TF-IDF Matrix",
        "4. 描述性统计",
        "5. LDA 建模与 K 值选择",
        "6. 主题解释与可视化",
    ], size=17)

    # 3
    slide = blank(prs)
    add_title(slide, "数据与预处理结果")
    data = [
        ["指标", "数值"],
        ["摘要数量", "12195"],
        ["词汇总数", "23818"],
        ["文档-词矩阵形状", "12195 × 23818"],
        ["非零元素数量", "878114"],
        ["稀疏度", "0.996977"],
        ["平均文档长度", "108.39"],
        ["文档长度中位数", "102"],
    ]
    add_table(slide, 0.85, 1.35, 5.0, 4.8, len(data), 2, data, font_size=14)
    add_bullets(slide, 6.35, 1.45, 5.95, 4.4, [
        "清洗阶段删除了 Background、Methods、Results、Conclusion 等结构化标题",
        "停用词表补充了通用英文停用词和部分领域高频词",
        "处理顺序调整为：词形还原后再次过滤停用词",
        "最终词汇表已去除 the、and、in、study、result 等无效高频词",
    ], size=17)

    # 4
    slide = blank(prs)
    add_title(slide, "预处理统计：高频词")
    add_image(slide, IMG["top_words"], 1.0, 1.25, w=11.3)
    add_textbox(slide, 0.95, 6.92, 11.4, 0.3, "高频词显示语料主要集中于 student、learning、teacher、program、teaching 等教育研究核心词。", size=13, color=MUTED)

    # 5
    slide = blank(prs)
    add_title(slide, "预处理统计：文档长度分布")
    add_image(slide, IMG["doc_len"], 1.0, 1.25, w=11.3)
    add_textbox(slide, 0.95, 6.92, 11.4, 0.3, "多数摘要预处理后长度集中在约 50-150 个词之间，整体适合进行文档级主题建模。", size=13, color=MUTED)

    # 6
    slide = blank(prs)
    add_title(slide, "LDA 主题模型")
    add_bullets(slide, 0.8, 1.35, 5.8, 5.3, [
        "使用 Python sklearn 的 LatentDirichletAllocation",
        "输入矩阵：Count Matrix",
        "输出每个主题的高概率词",
        "输出每篇摘要的主题分布",
        "提取每个主题的代表文档",
    ], size=18)
    data = [
        ["项目", "设置"],
        ["最终模型", "LDA"],
        ["最终主题数", "K = 15"],
        ["训练迭代", "max_iter = 50"],
        ["随机种子", "100"],
        ["输出目录", "output/lda_k15"],
    ]
    add_table(slide, 7.0, 1.45, 5.2, 3.6, len(data), 2, data, font_size=13)

    # 7
    slide = blank(prs)
    add_title(slide, "主题数量 K 的选择")
    add_image(slide, IMG["k_score"], 0.8, 1.3, w=5.9)
    add_image(slide, IMG["k_npmi"], 6.9, 1.3, w=5.7)
    add_textbox(slide, 0.9, 6.82, 11.7, 0.42, "K=1-4 作为粗粒度参考，不参与最终选择；在 K≥5 的候选模型中，K=15 综合评分最高。", size=13, color=MUTED)

    # 8
    slide = blank(prs)
    add_title(slide, "为什么选择 K=15")
    add_bullets(slide, 0.8, 1.25, 11.7, 4.8, [
        "K=15 在合格候选模型中取得最高综合分",
        "NPMI coherence 最高，说明主题内部关键词共现关系更好",
        "相比 K=5 和 K=10，K=15 提供更细致的主题划分",
        "相比 K=30、K=40、K=50，K=15 避免了过度碎片化",
        "最终综合考虑 coherence、exclusivity、topic separation、topic diversity 和 perplexity",
    ], size=20)
    add_textbox(slide, 0.8, 6.2, 11.6, 0.55, "结论：最终采用 LDA, K=15 作为主题模型。", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # 9
    slide = blank(prs)
    add_title(slide, "最终识别出的 15 个主题")
    topics_left = [
        "1 技术支持的在线与协作学习设计",
        "2 STEM 与工程科学教育",
        "3 教师教育、课程设计与专业发展",
        "4 化学与实验室科学教学",
        "5 在线学习中的自我调节、动机与参与",
        "6 学生成绩、测试与学习效果评估",
        "7 医学、护理与临床健康教育",
        "8 学业路径、职业发展与高等教育机构",
    ]
    topics_right = [
        "9 英语语言、写作与读写能力",
        "10 学生健康、性别与社会心理因素",
        "11 批判性思维、问题解决与身份发展",
        "12 同行反馈、学术出版与文献评审",
        "13 测量模型、量表验证与研究方法",
        "14 住院医师、外科与医学培训评价",
        "15 国际学生、中国语境与跨文化高等教育",
    ]
    add_bullets(slide, 0.65, 1.15, 6.1, 5.75, topics_left, size=15)
    add_bullets(slide, 6.95, 1.15, 5.9, 5.75, topics_right, size=15)

    # 10
    slide = blank(prs)
    add_title(slide, "每个主题的 Top Words")
    add_image(slide, IMG["topic_words"], 0.55, 1.05, w=12.2)

    # 11
    slide = blank(prs)
    add_title(slide, "主题平均占比与主导主题分布")
    add_image(slide, IMG["topic_prev"], 0.7, 1.3, w=5.9)
    add_image(slide, IMG["dominant"], 6.85, 1.3, w=5.9)
    add_textbox(slide, 0.9, 6.8, 11.5, 0.35, "Topic 3、Topic 1、Topic 15 等主题在语料中占比较高，是该摘要集合中的核心方向。", size=13, color=MUTED)

    # 12
    slide = blank(prs)
    add_title(slide, "主题相关性热力图")
    add_image(slide, IMG["corr"], 1.75, 1.05, h=6.0)
    add_textbox(slide, 8.85, 1.4, 3.7, 4.3, "热力图展示不同主题在文档中的共现关系。\n\n颜色越接近红色，说明两个主题在相同文档中共同出现的倾向越强；颜色越接近蓝色，说明二者更少共同出现。", size=15, color=TEXT)

    # 13
    slide = blank(prs)
    add_title(slide, "主题聚类树状图")
    add_image(slide, IMG["dendro"], 0.85, 1.2, w=11.7)
    add_textbox(slide, 0.95, 6.88, 11.3, 0.32, "树状图基于 1 - correlation 对主题进行层次聚类，用于观察哪些主题在文档层面更接近。", size=13, color=MUTED)

    # 14
    slide = blank(prs)
    add_title(slide, "文档-主题分布")
    add_image(slide, IMG["doc_topic"], 0.9, 1.05, w=11.6)
    add_textbox(slide, 0.95, 6.9, 11.4, 0.32, "每篇摘要不是只属于一个主题，而是由多个主题按比例混合组成。", size=13, color=MUTED)

    # 15
    slide = blank(prs)
    add_title(slide, "主要发现")
    add_bullets(slide, 0.85, 1.3, 11.6, 5.2, [
        "语料整体以教育研究为核心，围绕学生、学习、教师、教学、课程和评估展开",
        "技术支持学习、教师教育、国际高等教育和医学教育是较突出的主题方向",
        "LDA 主题模型能够将摘要划分为 15 个相对清晰的研究主题",
        "主题可视化显示部分主题之间存在明显关联，例如在线学习与学习动机、医学教育与住院医师培训",
    ], size=20)

    # 16
    slide = blank(prs)
    add_title(slide, "局限与后续工作")
    add_bullets(slide, 0.85, 1.25, 11.7, 5.4, [
        "当前数据主要包含摘要文本，缺少年份、期刊、国家等 metadata，因此无法做主题趋势分析",
        "仍有少量异常词，例如 engineere、reade、nurs、sic，可继续优化词形还原和停用词表",
        "如果课程要求使用 R，可以进一步用 STM 复现主题模型并输出 FREX words",
        "后续报告可结合主题解释表和可视化图进一步分析主题之间的关系",
    ], size=19)

    # 17
    slide = blank(prs)
    add_title(slide, "结论")
    add_bullets(slide, 0.95, 1.35, 11.3, 4.6, [
        "本项目完成了从原始英文摘要到主题解释的完整 NLP 流程",
        "预处理后得到 12195 篇摘要和 23818 个有效词汇",
        "通过多指标比较，最终选择 LDA K=15",
        "最终识别出 15 个教育研究相关主题，并完成主题命名、代表文档和可视化分析",
    ], size=21)
    add_textbox(slide, 1.4, 6.15, 10.5, 0.55, "Thank You", size=30, bold=True, color=TITLE, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"Saved presentation to {OUT}")


if __name__ == "__main__":
    make_ppt()
